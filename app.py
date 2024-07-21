# # import streamlit as st
# # import transformers
# # import torch
# # from dotenv import load_dotenv
# # from PyPDF2 import PdfReader
# # from langchain.text_splitter import CharacterTextSplitter
# # from langchain.embeddings import OpenAIEmbeddings, HuggingFaceInstructEmbeddings , HuggingFaceEmbeddings
# # from langchain.vectorstores import FAISS
# # from langchain.chat_models import ChatOpenAI
# # from langchain.memory import ConversationBufferMemory
# # from langchain.chains import ConversationalRetrievalChain
# # from htmlTemplates import css, bot_template, user_template
# # from langchain.llms import HuggingFaceHub

# # def get_pdf_text(pdf_docs):
# #     text = ""
# #     for pdf in pdf_docs:
# #         pdf_reader = PdfReader(pdf)
# #         for page in pdf_reader.pages:
# #             text += page.extract_text()
# #     return text


# # def get_text_chunks(text):
# #     text_splitter = CharacterTextSplitter(
# #         separator="\n",
# #         chunk_size=1000,
# #         chunk_overlap=200,
# #         length_function=len
# #     )
# #     chunks = text_splitter.split_text(text)
# #     return chunks


# # def get_vectorstore(text_chunks):
# #     embeddings = HuggingFaceEmbeddings()
# #     vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
# #     return vectorstore


# # def get_conversation_chain(vectorstore):
# #     retries = 3
# #     delay = 5  # Delay between retries in seconds

# #     for attempt in range(retries):
# #         try:
# #             llm = HuggingFaceHub(repo_id="google/flan-t5-xxl", model_kwargs={"temperature":0.5, "max_length":512})
# #             memory = ConversationBufferMemory(
# #                 memory_key='chat_history', return_messages=True)
# #             conversation_chain = ConversationalRetrievalChain.from_llm(
# #                 llm=llm,
# #                 retriever=vectorstore.as_retriever(),
# #                 memory=memory
# #             )
# #             return conversation_chain
# #         except Exception as e:
# #             if attempt < retries - 1:
# #                 time.sleep(delay)
# #             else:
# #                 raise e


# # def handle_userinput(user_question):
# #     response = st.session_state.conversation({'question': user_question})
# #     st.session_state.chat_history = response['chat_history']

# #     for i, message in enumerate(st.session_state.chat_history):
# #         if i % 2 == 0:
# #             st.write(user_template.replace(
# #                 "{{MSG}}", message.content), unsafe_allow_html=True)
# #         else:
# #             st.write(bot_template.replace(
# #                 "{{MSG}}", message.content), unsafe_allow_html=True)


# # def main():
# #     load_dotenv()
# #     st.set_page_config(page_title="Chat with multiple PDFs",
# #                        page_icon=":books:")
# #     st.write(css, unsafe_allow_html=True)

# #     if "conversation" not in st.session_state:
# #         st.session_state.conversation = None
# #     if "chat_history" not in st.session_state:
# #         st.session_state.chat_history = None

# #     st.header("Chat with multiple PDFs :books:")
# #     user_question = st.text_input("Ask a question about your documents:")
# #     if user_question:
# #         handle_userinput(user_question)

# #     with st.sidebar:
# #         st.subheader("Your documents")
# #         pdf_docs = st.file_uploader(
# #             "Upload your PDFs here and click on 'Process'", accept_multiple_files=True)
# #         if st.button("Process"):
# #             with st.spinner("Processing"):
# #                 # get pdf text
# #                 raw_text = get_pdf_text(pdf_docs)

# #                 # get the text chunks
# #                 text_chunks = get_text_chunks(raw_text)

# #                 # create vector store
# #                 vectorstore = get_vectorstore(text_chunks)

# #                 # create conversation chain
# #                 st.session_state.conversation = get_conversation_chain(
# #                     vectorstore)


# # if __name__ == '__main__':
# #     main()

# # import streamlit as st
# # import transformers
# # import torch
# # from dotenv import load_dotenv
# # from PyPDF2 import PdfReader
# # from langchain.text_splitter import CharacterTextSplitter
# # from langchain.embeddings import HuggingFaceEmbeddings
# # from langchain.vectorstores import FAISS
# # from langchain.chat_models import ChatOpenAI
# # from langchain.memory import ConversationBufferMemory
# # from langchain.chains import ConversationalRetrievalChain
# # from htmlTemplates import css, bot_template, user_template
# # from langchain.llms import HuggingFaceHub

# # def get_pdf_text(pdf_docs):
# #     text = ""
# #     for pdf in pdf_docs:
# #         pdf_reader = PdfReader(pdf)
# #         for page in pdf_reader.pages:
# #             text += page.extract_text()
# #     return text

# # def get_text_chunks(text):
# #     text_splitter = CharacterTextSplitter(
# #         separator="\n",
# #         chunk_size=1000,
# #         chunk_overlap=200,
# #         length_function=len
# #     )
# #     chunks = text_splitter.split_text(text)
# #     return chunks

# # def get_vectorstore(text_chunks):
# #     embeddings = HuggingFaceEmbeddings()
# #     vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
# #     return vectorstore

# # def get_conversation_chain(vectorstore):
# #     # llm = ChatOpenAI()
# #     llm = HuggingFaceHub(repo_id="google/flan-t5-xxl", model_kwargs={"temperature":0.5, "max_length":512})

# #     memory = ConversationBufferMemory(
# #         memory_key='chat_history', return_messages=True)
# #     conversation_chain = ConversationalRetrievalChain.from_llm(
# #         llm=llm,
# #         retriever=vectorstore.as_retriever(),
# #         memory=memory
# #     )
# #     return conversation_chain

# # def handle_userinput(user_question):
# #     if st.session_state.conversation is not None:
# #         response = st.session_state.conversation({'question': user_question})
# #         st.session_state.chat_history = response['chat_history']

# #         for i, message in enumerate(st.session_state.chat_history):
# #             if i % 2 == 0:
# #                 st.write(user_template.replace(
# #                     "{{MSG}}", message.content), unsafe_allow_html=True)
# #             else:
# #                 st.write(bot_template.replace(
# #                     "{{MSG}}", message.content), unsafe_allow_html=True)
# #     else:
# #         st.error("Conversation chain is not initialized. Please upload and process the PDFs first.")

# # def main():
# #     load_dotenv()
# #     st.set_page_config(page_title="Chat with multiple PDFs", page_icon=":books:")
# #     st.write(css, unsafe_allow_html=True)

# #     if "conversation" not in st.session_state:
# #         st.session_state.conversation = None
# #     if "chat_history" not in st.session_state:
# #         st.session_state.chat_history = None

# #     st.header("Chat with multiple PDFs :books:")
# #     user_question = st.text_input("Ask a question about your documents:")
# #     if user_question:
# #         handle_userinput(user_question)

# #     with st.sidebar:
# #         st.subheader("Your documents")
# #         pdf_docs = st.file_uploader(
# #             "Upload your PDFs here and click on 'Process'", accept_multiple_files=True)
# #         if st.button("Process"):
# #             if pdf_docs:
# #                 with st.spinner("Processing"):
# #                     # Get pdf text
# #                     raw_text = get_pdf_text(pdf_docs)

# #                     # Get the text chunks
# #                     text_chunks = get_text_chunks(raw_text)

# #                     # Create vector store
# #                     vectorstore = get_vectorstore(text_chunks)

# #                     # Create conversation chain
# #                     st.session_state.conversation = get_conversation_chain(vectorstore)
# #                     st.success("PDFs processed and conversation chain initialized.")
# #             else:
# #                 st.error("Please upload at least one PDF document.")

# # if __name__ == '__main__':
# #     main()

# import time
# import streamlit as st
# from dotenv import load_dotenv
# from PyPDF2 import PdfReader
# from langchain.text_splitter import CharacterTextSplitter
# from langchain.embeddings import HuggingFaceEmbeddings
# from langchain.vectorstores import FAISS
# from langchain.chat_models import ChatOpenAI
# from langchain.memory import ConversationBufferMemory
# from langchain.chains import ConversationalRetrievalChain
# from htmlTemplates import css, bot_template, user_template
# from langchain.llms import HuggingFaceHub

# def get_pdf_text(pdf_docs):
#     text = ""
#     for pdf in pdf_docs:
#         pdf_reader = PdfReader(pdf)
#         for page in pdf_reader.pages:
#             text += page.extract_text()
#     return text

# def get_text_chunks(text):
#     text_splitter = CharacterTextSplitter(
#         separator="\n",
#         chunk_size=1000,
#         chunk_overlap=200,
#         length_function=len
#     )
#     chunks = text_splitter.split_text(text)
#     return chunks

# def get_vectorstore(text_chunks):
#     embeddings = HuggingFaceEmbeddings()
#     vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
#     return vectorstore

# def get_conversation_chain(vectorstore):
#     retries = 3
#     delay = 5  # Delay between retries in seconds

#     for attempt in range(retries):
#         try:
#             llm = HuggingFaceHub(repo_id="google/flan-t5-large", model_kwargs={"temperature":0.5, "max_length":512})
#             memory = ConversationBufferMemory(
#                 memory_key='chat_history', return_messages=True)
#             conversation_chain = ConversationalRetrievalChain.from_llm(
#                 llm=llm,
#                 retriever=vectorstore.as_retriever(),
#                 memory=memory
#             )
#             return conversation_chain
#         except Exception as e:
#             if attempt < retries - 1:
#                 time.sleep(delay)
#             else:
#                 st.error(f"Failed to initialize the conversation chain after {retries} attempts. Please try again later.")
#                 raise e

# def handle_userinput(user_question):
#     if st.session_state.conversation is not None:
#         try:
#             response = st.session_state.conversation({'question': user_question})
#             st.session_state.chat_history = response['chat_history']

#             for i, message in enumerate(st.session_state.chat_history):
#                 if i % 2 == 0:
#                     st.write(user_template.replace(
#                         "{{MSG}}", message.content), unsafe_allow_html=True)
#                 else:
#                     st.write(bot_template.replace(
#                         "{{MSG}}", message.content), unsafe_allow_html=True)
#         except Exception as e:
#             st.error("An error occurred while processing your request.")
#             st.error(f"Error details: {e}")
#     else:
#         st.error("Conversation chain is not initialized. Please upload and process the PDFs first.")

# def main():
#     load_dotenv()
#     st.set_page_config(page_title="Chat with multiple PDFs", page_icon=":books:")
#     st.write(css, unsafe_allow_html=True)

#     if "conversation" not in st.session_state:
#         st.session_state.conversation = None
#     if "chat_history" not in st.session_state:
#         st.session_state.chat_history = None

#     st.header("Chat with multiple PDFs :books:")
#     user_question = st.text_input("Ask a question about your documents:")
#     if user_question:
#         handle_userinput(user_question)

#     with st.sidebar:
#         st.subheader("Your documents")
#         pdf_docs = st.file_uploader(
#             "Upload your PDFs here and click on 'Process'", accept_multiple_files=True)
#         if st.button("Process"):
#             if pdf_docs:
#                 with st.spinner("Processing"):
#                     # Get pdf text
#                     raw_text = get_pdf_text(pdf_docs)

#                     # Get the text chunks
#                     text_chunks = get_text_chunks(raw_text)

#                     # Create vector store
#                     vectorstore = get_vectorstore(text_chunks)

#                     # Create conversation chain
#                     st.session_state.conversation = get_conversation_chain(vectorstore)
#                     st.success("PDFs processed and conversation chain initialized.")
#             else:
#                 st.error("Please upload at least one PDF document.")

# if __name__ == '__main__':
#     main()

import streamlit as st
from htmlTemplates import css, bot_template, user_template
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_huggingface.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain.chains import create_retrieval_chain
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain_core.prompts import ChatPromptTemplate
import os
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
import tempfile

model_name = "sentence-transformers/all-MiniLM-L6-v2"
model_kwargs = {'device': 'cpu'}
encode_kwargs = {'normalize_embeddings': False}
embeddings = HuggingFaceEmbeddings(
    model_name=model_name,
    model_kwargs=model_kwargs,
    encode_kwargs=encode_kwargs
)

def get_pdf_text(pdf_file):
    text = ""
    pdf_reader = PdfReader(pdf_file)
    for page_num, page in enumerate(pdf_reader.pages, 1):
        text += f"Page {page_num}:\n"
        text += page.extract_text()
        text += "\n\n"  # Add space between pages
    return text

def extract_text_from_pptx(pptx_file):
    presentation = Presentation(pptx_file)
    text = ""
    for slide_num, slide in enumerate(presentation.slides, 1):
        text += f"Slide {slide_num}:\n"
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text += f"{paragraph.text}\n"
        text += "\n"  # Add a separator between slides
    return text

def extract_text_from_docx(docx_file):
    doc = Document(docx_file)
    text = ""
    for para in doc.paragraphs:
        if para.style.name.startswith('Heading'):
            text += f"\n\n{para.text}\n\n"  # Add extra space around headings
        else:
            text += para.text + " "  # Combine lines into paragraphs

    # Ensure paragraphs are separated properly
    paragraphs = text.split("\n\n")
    formatted_text = "\n\n".join(para.strip() for para in paragraphs if para.strip())

    for table in doc.tables:
        formatted_text += "\n\n"
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    formatted_text += cell.text + "\t"
            formatted_text += "\n"
        formatted_text += "\n"

    return formatted_text

def get_response(question):
    db = FAISS.load_local("vectorstore", embeddings, allow_dangerous_deserialization=True)
    system_prompt = (
        "Use the given context to answer the question. "
        "If you don't know the answer, say you I cant answer the question right now as I didnt understand the question. "
        "Use five sentence maximum and keep the answer detailed. "
        "Context: {context}"
    )
    prompt = ChatPromptTemplate.from_messages([
        ("system", system_prompt),
        ("human", "{input}"),
    ])

    llm = ChatGoogleGenerativeAI(
        model="gemini-1.5-pro",
        temperature=0.5,
        max_tokens=None,
        timeout=None,
        max_retries=2,
        google_api_key="AIzaSyByAULS9YrPUqkrai_ZQn5-PPaOxBaTpcU"
    )

    retriever = db.as_retriever()
    question_answer_chain = create_stuff_documents_chain(llm, prompt)
    chain = create_retrieval_chain(retriever, question_answer_chain)

    response = chain.invoke({"input": question})
    return response['answer']

def handle_userinput(user_question):
    response = get_response(user_question)
    st.session_state.chat_history.append({"role": "user", "content": user_question})
    st.session_state.chat_history.append({"role": "assistant", "content": response})

    for message in st.session_state.chat_history:
        if message["role"] == "user":
            st.write(user_template.replace("{{MSG}}", message["content"]), unsafe_allow_html=True)
        else:
            st.write(bot_template.replace("{{MSG}}", message["content"]), unsafe_allow_html=True)

def main():
    st.set_page_config(page_title="Convo AI", page_icon=":books:")
    st.write(css, unsafe_allow_html=True)

    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []

    st.header("Welcome to Convo AI , Let's Chat !!")

    user_question = st.text_input("Ask a question about your documents:")
    if user_question:
        handle_userinput(user_question)

    with st.sidebar:
        st.subheader("Your documents")
        uploaded_files = st.file_uploader("Upload your PDFs, PPTs, or Word files here", accept_multiple_files=True, type=['pdf', 'pptx', 'docx'])
        
        if st.button("Process"):
            if uploaded_files:
                for file in uploaded_files:
                    with tempfile.NamedTemporaryFile(delete=False, suffix=file.name[-5:]) as temp_file:
                        temp_file.write(file.getvalue())
                        temp_file_path = temp_file.name

                    text_final = ""
                    
                    if file.name.endswith('.pdf'):
                        text = get_pdf_text(temp_file_path).strip()
                    elif file.name.endswith('.pptx'):
                        text = extract_text_from_pptx(temp_file_path).strip()
                    elif file.name.endswith('.docx'):
                        text = extract_text_from_docx(temp_file_path).strip()
                    else:
                        text = ""

                    text_final += text
                    os.unlink(temp_file_path)  # Delete the temporary file

                text_splitter = RecursiveCharacterTextSplitter(chunk_size=900, chunk_overlap=100)
                texts = text_splitter.split_text(text_final)

                db = FAISS.from_texts(texts, embeddings)
                db.save_local("vectorstore")

                st.success("Documents processed successfully!")
            else:
                st.warning("Please upload at least one file.")

if __name__ == "__main__":
    main()